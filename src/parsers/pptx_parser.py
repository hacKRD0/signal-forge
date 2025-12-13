"""PPTX document parser.

This module provides functionality to extract text content from PowerPoint
PPTX files using python-pptx library.
"""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.exc import PackageNotFoundError


logger = logging.getLogger(__name__)


def parse_pptx(file_path: str) -> str:
    """Extract text content from a PowerPoint PPTX file.

    Reads a PowerPoint PPTX file and extracts all text content from slides,
    including text from shapes and notes. Maintains slide structure with
    clear slide separators.

    Args:
        file_path: Path to the PPTX file to parse (absolute or relative).

    Returns:
        Extracted text as a string with slide separators ("Slide N:").
        Returns empty string if the presentation has no text.

    Raises:
        FileNotFoundError: If the specified file does not exist.
        ValueError: If the file is corrupted or not a valid PPTX file.
        PermissionError: If the file cannot be read due to permissions.

    Example:
        >>> text = parse_pptx('presentations/pitch.pptx')
        >>> print(f"Extracted {len(text)} characters")

    Note:
        This parser extracts text from all shapes in each slide and preserves
        the slide-by-slide structure for context.
    """
    # Convert to Path for better path handling
    path = Path(file_path)

    # Check if file exists
    if not path.exists():
        logger.error(f"File not found: {file_path}")
        raise FileNotFoundError(f"File not found: {file_path}")

    logger.info(f"Parsing PPTX file: {file_path}")

    try:
        # Open the presentation
        prs = Presentation(str(path))

        # Get number of slides
        num_slides = len(prs.slides)
        logger.debug(f"Presentation has {num_slides} slides")

        if num_slides == 0:
            logger.warning(f"Presentation has no slides: {file_path}")
            return ""

        # Extract text from all slides
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []

            # Add slide separator
            slide_text.append(f"--- Slide {slide_num} ---")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            # Extract notes if present
            if slide.has_notes_slide:
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_text.append(f"[Notes: {notes_text}]")
                except Exception as e:
                    logger.debug(f"Could not extract notes from slide {slide_num}: {e}")

            # Only add slide if it has content
            if len(slide_text) > 1:  # More than just the separator
                text_parts.append("\n".join(slide_text))
            else:
                logger.debug(f"Slide {slide_num} has no extractable text")

        # Join all slide parts with double newlines
        extracted_text = "\n\n".join(text_parts)

        logger.info(
            f"Successfully extracted {len(extracted_text)} characters "
            f"from {num_slides} slides in {file_path}"
        )

        return extracted_text

    except PackageNotFoundError as e:
        logger.error(f"Corrupted or invalid PPTX file: {file_path}")
        raise ValueError(
            f"File is corrupted or not a valid PPTX file: {file_path}"
        ) from e
    except PermissionError as e:
        logger.error(f"Permission denied reading file: {file_path}")
        raise
    except Exception as e:
        logger.error(f"Unexpected error parsing PPTX file {file_path}: {e}")
        raise ValueError(
            f"Error parsing PPTX file {file_path}: {str(e)}"
        ) from e
